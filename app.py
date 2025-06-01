from flask import Flask, render_template, request, send_file, redirect, url_for, flash, get_flashed_messages
import os
import sys
from datetime import datetime
import zipfile
import logging
import pandas as pd
from werkzeug.utils import secure_filename

# Conversion specific imports
from pdf2docx import Converter as PdfToDocxConverter
import pdfplumber
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from pptx import Presentation

# Configure Poppler path
POPPLER_PATH = r'C:\Users\bke\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin'  # Update this to your actual Poppler path
if not os.path.exists(POPPLER_PATH):
    logging.warning("Poppler path not found. Image conversions may not work.")

try:
    import pypandoc
    PANDOC_AVAILABLE = True
except (ImportError, OSError):
    PANDOC_AVAILABLE = False
    logging.warning("Pypandoc or Pandoc not available. Some conversions will be disabled.")

# Tabula for table extraction
try:
    import tabula
    TABULA_AVAILABLE = True
except ImportError:
    TABULA_AVAILABLE = False
    logging.warning("tabula-py not available. CSV/XLSX conversion will be disabled.")

app = Flask(__name__)
app.secret_key = os.urandom(24)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Configuration
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['CONVERTED_FOLDER'] = 'converted'
app.config['ALLOWED_EXTENSIONS'] = {'pdf'}
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB limit
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['CONVERTED_FOLDER'], exist_ok=True)

# Supported Formats
SUPPORTED_FORMATS = {
    "docx": {"name": "Word (.docx)", "ext": "docx", "status": "stable"},
    "doc": {"name": "Word 97-2003 (.doc)", "ext": "doc",
            "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
    "txt": {"name": "Plain Text (.txt)", "ext": "txt", "status": "stable"},
    "rtf": {"name": "Rich Text Format (.rtf)", "ext": "rtf",
            "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
    "epub": {"name": "EPUB (.epub)", "ext": "epub", "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
    "mobi": {"name": "MOBI (.mobi)", "ext": "mobi", "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
    "html": {"name": "HTML (.html)", "ext": "html", "status": "stable"},
    "xml": {"name": "XML (.xml)", "ext": "xml", "status": "stable"},
    "csv": {"name": "CSV (.csv)", "ext": "csv", "status": "requires_java" if TABULA_AVAILABLE else "disabled"},
    "xlsx": {"name": "Excel (.xlsx)", "ext": "xlsx", "status": "requires_java" if TABULA_AVAILABLE else "disabled"},
    "pptx": {"name": "PowerPoint (.pptx)", "ext": "pptx", "status": "stable"},
    "ppt": {"name": "PowerPoint 97-2003 (.ppt)", "ext": "ppt",
            "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
    "png": {"name": "PNG Images (.zip)", "ext": "zip", "status": "stable"},
    "jpeg": {"name": "JPEG Images (.zip)", "ext": "zip", "status": "stable"},
    "tiff": {"name": "TIFF Images (.zip)", "ext": "zip", "status": "stable"},
    "svg": {"name": "SVG Images (.zip)", "ext": "zip", "status": "stable"},
    "odt": {"name": "OpenDocument Text (.odt)", "ext": "odt",
            "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
    "markdown": {"name": "Markdown (.md)", "ext": "md",
                 "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
    "pages": {"name": "Apple Pages (.pages)", "ext": "pages",
              "status": "requires_pandoc" if PANDOC_AVAILABLE else "disabled"},
}

def allowed_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def generate_html_content(title, body_content):
    return f"""<!DOCTYPE html>
<html lang="en"><head><meta charset="UTF-8"><title>{title}</title></head>
<body>{body_content}</body></html>"""

def get_safe_filename_base(filename):
    base = os.path.splitext(filename)[0]
    safe_base = "".join(c if c.isalnum() or c in [' ', '_', '-'] else '_' for c in base).strip()
    return safe_base[:50] if safe_base else "file"

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        if 'pdf_file' not in request.files:
            flash('No file selected', 'error')
            return redirect(request.url)

        file = request.files['pdf_file']
        output_format = request.form.get('output_format')

        if file.filename == '':
            flash('No file selected', 'error')
            return redirect(request.url)

        if not output_format or output_format not in SUPPORTED_FORMATS:
            flash('Invalid output format', 'error')
            return redirect(request.url)

        if not allowed_file(file.filename):
            flash('Only PDF files are allowed', 'error')
            return redirect(request.url)

        # Process the file
        filename = secure_filename(file.filename)
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S%f")
        safe_name = get_safe_filename_base(filename)

        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], f"upload_{timestamp}_{safe_name}.pdf")
        file.save(pdf_path)

        output_ext = SUPPORTED_FORMATS[output_format]['ext']
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], f"converted_{timestamp}_{safe_name}.{output_ext}")

        try:
            if output_format == 'docx':
                cv = PdfToDocxConverter(pdf_path)
                cv.convert(output_path)
                cv.close()

            elif output_format == 'txt':
                with pdfplumber.open(pdf_path) as pdf:
                    text = "\n\n".join(page.extract_text() or "" for page in pdf.pages)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)

            elif output_format == 'html':
                doc = fitz.open(pdf_path)
                html = "".join(f"<div>{page.get_text('html')}</div>" for page in doc)
                doc.close()
                full_html = generate_html_content(safe_name, html)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(full_html)

            elif output_format == 'xml':
                doc = fitz.open(pdf_path)
                xml = doc.get_text("xml")
                doc.close()
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(xml)

            elif output_format in ['png', 'jpeg', 'tiff']:
                fmt = 'png' if output_format == 'png' else 'jpeg' if output_format == 'jpeg' else 'tiff'
                images = convert_from_path(pdf_path, fmt=fmt, poppler_path=POPPLER_PATH)
                temp_dir = f"temp_{timestamp}"
                os.makedirs(temp_dir, exist_ok=True)

                for i, img in enumerate(images):
                    img.save(os.path.join(temp_dir, f"page_{i + 1}.{fmt}"), fmt.upper())

                with zipfile.ZipFile(output_path, 'w') as zipf:
                    for root, _, files in os.walk(temp_dir):
                        for file in files:
                            zipf.write(os.path.join(root, file), file)

                for file in os.listdir(temp_dir):
                    os.remove(os.path.join(temp_dir, file))
                os.rmdir(temp_dir)

            elif output_format == 'svg':
                doc = fitz.open(pdf_path)
                temp_dir = f"temp_{timestamp}"
                os.makedirs(temp_dir, exist_ok=True)

                for i, page in enumerate(doc):
                    svg = page.get_svg_image()
                    with open(os.path.join(temp_dir, f"page_{i + 1}.svg"), 'w') as f:
                        f.write(svg)

                with zipfile.ZipFile(output_path, 'w') as zipf:
                    for file in os.listdir(temp_dir):
                        zipf.write(os.path.join(temp_dir, file), file)

                for file in os.listdir(temp_dir):
                    os.remove(os.path.join(temp_dir, file))
                os.rmdir(temp_dir)
                doc.close()

            elif output_format == 'pptx':
                prs = Presentation()
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout

                images = convert_from_path(pdf_path, poppler_path=POPPLER_PATH)
                temp_dir = f"temp_{timestamp}"
                os.makedirs(temp_dir, exist_ok=True)

                for i, img in enumerate(images):
                    img_path = os.path.join(temp_dir, f"page_{i + 1}.png")
                    img.save(img_path, 'PNG')

                    slide = prs.slides.add_slide(blank_slide_layout)
                    left = top = 0
                    pic = slide.shapes.add_picture(img_path, left, top,
                                                   width=prs.slide_width,
                                                   height=prs.slide_height)

                prs.save(output_path)

                for file in os.listdir(temp_dir):
                    os.remove(os.path.join(temp_dir, file))
                os.rmdir(temp_dir)

            elif PANDOC_AVAILABLE and output_format in ['doc', 'rtf', 'epub', 'mobi', 'odt', 'markdown', 'pages']:
                doc = fitz.open(pdf_path)
                text = "\n\n".join(page.get_text() for page in doc)
                doc.close()

                if output_format == 'markdown':
                    format = 'gfm'
                elif output_format == 'pages':
                    format = 'pages'
                else:
                    format = output_format

                pypandoc.convert_text(text, format, format='plain', outputfile=output_path)

            elif TABULA_AVAILABLE and output_format == 'csv':
                dfs = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
                if len(dfs) == 1:
                    dfs[0].to_csv(output_path, index=False)
                else:
                    temp_dir = f"temp_{timestamp}"
                    os.makedirs(temp_dir, exist_ok=True)

                    for i, df in enumerate(dfs):
                        df.to_csv(os.path.join(temp_dir, f"table_{i + 1}.csv"), index=False)

                    with zipfile.ZipFile(output_path, 'w') as zipf:
                        for file in os.listdir(temp_dir):
                            zipf.write(os.path.join(temp_dir, file), file)

                    for file in os.listdir(temp_dir):
                        os.remove(os.path.join(temp_dir, file))
                    os.rmdir(temp_dir)

            elif TABULA_AVAILABLE and output_format == 'xlsx':
                dfs = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
                with pd.ExcelWriter(output_path) as writer:
                    for i, df in enumerate(dfs):
                        df.to_excel(writer, sheet_name=f"Table_{i + 1}", index=False)

            flash('Conversion successful!', 'success')
            return render_template('index.html',
                                   download_file=f"converted_{timestamp}_{safe_name}.{output_ext}",
                                   supported_formats=SUPPORTED_FORMATS,
                                   pandoc_available=PANDOC_AVAILABLE,
                                   tabula_available=TABULA_AVAILABLE)

        except Exception as e:
            logging.error(f"Conversion error: {str(e)}", exc_info=True)
            flash(f'Conversion failed: {str(e)}', 'error')
            if os.path.exists(output_path):
                os.remove(output_path)

    return render_template('index.html',
                           supported_formats=SUPPORTED_FORMATS,
                           pandoc_available=PANDOC_AVAILABLE,
                           tabula_available=TABULA_AVAILABLE)

@app.route('/download/<filename>')
def download(filename):
    safe_filename = secure_filename(filename)
    if not safe_filename:
        flash('Invalid filename', 'error')
        return redirect(url_for('index'))

    file_path = os.path.join(app.config['CONVERTED_FOLDER'], safe_filename)
    if not os.path.exists(file_path):
        flash('File not found', 'error')
        return redirect(url_for('index'))

    return send_file(file_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)